import os
import pptx
import argparse
from page_eval import eval_page, viz_scores

def main():
    ref_path = [f for f in os.listdir(args.example_dir) if f.endswith(".pptx")]
    assert len(ref_path) == 1, "There should be exactly one PPTX file in the example directory."
    ref_path = os.path.join(args.example_dir, ref_path[0])

    page_list = [
        p for p in os.listdir(args.example_dir) 
        if ('.' not in p) and (int(p.split('_')[-1]) < args.max_pages)
    ]
    page_list = [os.path.join(args.example_dir, p, f"{args.predict_name}.pptx") for p in page_list]
    page_list = [p for p in page_list if os.path.exists(p)]
    ref_prs = pptx.Presentation(ref_path)
    print(f"Execution success rate: {100*len(page_list)/min(args.max_pages, len(ref_prs.slides)):4.1f}%")
    
    page_list = sorted(page_list, key=lambda x: int(x.split('/')[-2].split('_')[1]))

    ref_prs = pptx.Presentation(ref_path)
    scores = {"match": [], "text": [], "color": [], "position": []}
    for p in page_list:
        idx = int(p.split('/')[-2].split('_')[1]) - 1
        gen_prs = pptx.Presentation(p)
        p_scores = eval_page(gen_prs, gen_page=0, ref_prs=ref_prs, ref_page=idx)
        for k, v in p_scores.items():
            scores[k].append(v)
    viz_scores(scores)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Evaluate PPTX file.")
    parser.add_argument("--example_dir", type=str, required=True)
    parser.add_argument("--predict_name", type=str, default="gpt-4o-mini")
    parser.add_argument("--max_pages", type=int, default=100)

    args = parser.parse_args()

    main()
